"""
File conversion engine for Rikkahub.

Supported conversions:
  txt ↔ md ↔ html (native)
  txt → docx, md → docx, html → docx
  pdf → txt, pdf → md, pdf → docx
  docx → txt, docx → md
  xlsx ↔ csv, xlsx → json, csv → json, json → csv, json → xlsx
  pptx → txt, pptx → md
  zip → extract
  svg → png, svg → jpg
  png ↔ jpg ↔ webp (image conversions)
  txt → pdf, md → pdf, html → pdf (via fpdf2)
  image → pdf (multiple images to one PDF)
  epub → txt, epub → md

Returns: {'stdout': str, 'files': [str], 'error': str?}
"""

import json
import sys
import os
import csv
import re
import base64
import zipfile
import traceback


def convert(input_path, input_text, from_format, to_format, output_dir):
    """File conversion entry point. Returns {'stdout': str, 'files': [str]}"""
    result = {'stdout': '', 'files': [], 'error': None}

    try:
        # ── Text-native conversions (Kotlin handles txt/md/html interop) ──
        # Python handles: anything involving docx, pdf, xlsx, pptx, zip, svg, epub, image→pdf

        # ── DOCX conversions ──
        if from_format == 'txt' and to_format == 'docx':
            from docx import Document
            text = _read_file(input_path, input_text)
            doc = Document()
            for para in text.split('\n\n'):
                p = para.strip()
                if not p:
                    continue
                if p.startswith('# '):
                    doc.add_heading(p[2:], level=1)
                elif p.startswith('## '):
                    doc.add_heading(p[3:], level=2)
                elif p.startswith('### '):
                    doc.add_heading(p[4:], level=3)
                else:
                    doc.add_paragraph(p)
            out = _outpath(input_path, 'docx', output_dir)
            doc.save(out)
            result['files'].append(out)
            result['stdout'] = f'Saved: {out}'

        elif from_format == 'md' and to_format == 'docx':
            from docx import Document
            text = _read_file(input_path, input_text)
            doc = Document()
            for line in text.split('\n'):
                s = line.strip()
                if s.startswith('# '):
                    doc.add_heading(s[2:], level=1)
                elif s.startswith('## '):
                    doc.add_heading(s[3:], level=2)
                elif s.startswith('### '):
                    doc.add_heading(s[4:], level=3)
                elif s.startswith('- ') or s.startswith('* '):
                    doc.add_paragraph(s, style='List Bullet')
                else:
                    doc.add_paragraph(s)
            out = _outpath(input_path, 'docx', output_dir)
            doc.save(out)
            result['files'].append(out)
            result['stdout'] = f'Saved: {out}'

        elif from_format == 'html' and to_format == 'docx':
            from docx import Document
            from bs4 import BeautifulSoup
            text = _read_file(input_path, input_text)
            soup = BeautifulSoup(text, 'html.parser')
            doc = Document()
            for el in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li']):
                tag = el.name
                txt = el.get_text(strip=True)
                if not txt:
                    continue
                if tag in ('h1',):
                    doc.add_heading(txt, level=1)
                elif tag in ('h2',):
                    doc.add_heading(txt, level=2)
                elif tag in ('h3', 'h4', 'h5', 'h6'):
                    doc.add_heading(txt, level=int(tag[1]))
                elif tag == 'li':
                    doc.add_paragraph(txt, style='List Bullet')
                else:
                    doc.add_paragraph(txt)
            out = _outpath(input_path, 'docx', output_dir)
            doc.save(out)
            result['files'].append(out)
            result['stdout'] = f'Saved: {out}'

        # ── PDF extraction ──
        elif from_format == 'pdf' and to_format in ('txt', 'md'):
            from pypdf import PdfReader
            reader = PdfReader(input_path)
            text = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text.append(t)
            output = '\n\n'.join(text)
            if to_format == 'md':
                output = '# Extracted from PDF\n\n' + output
            result['stdout'] = output

        elif from_format == 'pdf' and to_format == 'docx':
            # PDF → DOCX: extract text, create docx
            from pypdf import PdfReader
            from docx import Document
            reader = PdfReader(input_path)
            doc = Document()
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    for para in text.split('\n'):
                        p = para.strip()
                        if p:
                            doc.add_paragraph(p)
            out = _outpath(input_path, 'docx', output_dir)
            doc.save(out)
            result['files'].append(out)
            result['stdout'] = f'Saved: {out}'

        # ── DOCX extraction ──
        elif from_format == 'docx' and to_format in ('txt', 'md'):
            from docx import Document
            doc = Document(input_path)
            lines = []
            for p in doc.paragraphs:
                t = p.text.strip()
                if not t:
                    continue
                if to_format == 'md':
                    style = p.style.name.lower() if p.style else ''
                    if 'heading 1' in style:
                        lines.append(f'# {t}')
                    elif 'heading 2' in style:
                        lines.append(f'## {t}')
                    elif 'heading 3' in style:
                        lines.append(f'### {t}')
                    elif 'list' in style:
                        lines.append(f'- {t}')
                    else:
                        lines.append(t)
                else:
                    lines.append(t)
            result['stdout'] = '\n'.join(lines)

        # ── Spreadsheet conversions ──
        elif from_format == 'xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(input_path)
            ws = wb.active
            if to_format == 'csv':
                out = _outpath(input_path, 'csv', output_dir)
                with open(out, 'w', newline='', encoding='utf-8') as f:
                    w = csv.writer(f)
                    for row in ws.iter_rows(values_only=True):
                        w.writerow(row)
                result['files'].append(out)
                result['stdout'] = f'Saved: {out}'
            elif to_format == 'json':
                headers = [c.value for c in ws[1]]
                data = []
                for row in ws.iter_rows(min_row=2, values_only=True):
                    data.append({headers[i]: row[i] for i in range(len(headers)) if i < len(headers)})
                result['stdout'] = json.dumps(data, ensure_ascii=False, indent=2)
            else:
                raise ValueError(f'xlsx → {to_format} not supported')

        elif from_format == 'csv':
            if to_format == 'json':
                with open(input_path, 'r', encoding='utf-8') as f:
                    reader = csv.DictReader(f)
                    data = list(reader)
                result['stdout'] = json.dumps(data, ensure_ascii=False, indent=2)
            elif to_format == 'xlsx':
                import openpyxl
                wb = openpyxl.Workbook()
                ws = wb.active
                with open(input_path, 'r', encoding='utf-8') as f:
                    reader = csv.reader(f)
                    for row in reader:
                        ws.append(row)
                out = _outpath(input_path, 'xlsx', output_dir)
                wb.save(out)
                result['files'].append(out)
                result['stdout'] = f'Saved: {out}'
            else:
                raise ValueError(f'csv → {to_format} not supported')

        elif from_format == 'json':
            data = json.load(open(input_path, 'r', encoding='utf-8'))
            if isinstance(data, dict):
                data = [data]
            if not data:
                raise ValueError('Empty JSON data')
            if to_format == 'csv':
                headers = list(data[0].keys())
                out = _outpath(input_path, 'csv', output_dir)
                with open(out, 'w', newline='', encoding='utf-8') as f:
                    w = csv.writer(f)
                    w.writerow(headers)
                    for row in data:
                        w.writerow([row.get(h, '') for h in headers])
                result['files'].append(out)
                result['stdout'] = f'Saved: {out}'
            elif to_format == 'xlsx':
                import openpyxl
                wb = openpyxl.Workbook()
                ws = wb.active
                headers = list(data[0].keys())
                ws.append(headers)
                for row in data:
                    ws.append([row.get(h, '') for h in headers])
                out = _outpath(input_path, 'xlsx', output_dir)
                wb.save(out)
                result['files'].append(out)
                result['stdout'] = f'Saved: {out}'
            else:
                raise ValueError(f'json → {to_format} not supported')

        # ── PowerPoint extraction ──
        elif from_format == 'pptx' and to_format in ('txt', 'md'):
            from pptx import Presentation
            prs = Presentation(input_path)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                if to_format == 'md':
                    lines.append(f'## Slide {i}')
                else:
                    lines.append(f'--- Slide {i} ---')
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        lines.append(shape.text.strip())
            result['stdout'] = '\n\n'.join(lines)

        # ── ZIP extraction ──
        elif from_format == 'zip':
            extract_dir = os.path.join(
                output_dir, os.path.basename(input_path).rsplit('.', 1)[0] + '_extracted'
            )
            with zipfile.ZipFile(input_path, 'r') as z:
                z.extractall(extract_dir)
            files = []
            for root, dirs, fnames in os.walk(extract_dir):
                for fname in fnames:
                    files.append(os.path.join(root, fname))
            result['stdout'] = f'Extracted to {extract_dir} ({len(files)} files)'
            result['files'] = files

        # ── SVG → PNG/JPG ──
        elif from_format == 'svg' and to_format in ('png', 'jpg'):
            # Use cairosvg if available, otherwise fall back to manual approach
            try:
                import cairosvg
            except ImportError:
                # Fallback: report available approach
                raise ValueError(
                    'SVG conversion requires the cairosvg library. '
                    'Install it or use the image conversion tool instead.'
                )
            svg_data = open(input_path, 'rb').read()
            out = _outpath(input_path, to_format, output_dir)
            if to_format == 'png':
                cairosvg.svg2png(bytestring=svg_data, write_to=out)
            else:
                import io
                png_data = cairosvg.svg2png(bytestring=svg_data)
                from PIL import Image
                img = Image.open(io.BytesIO(png_data))
                img = img.convert('RGB')
                img.save(out, 'JPEG', quality=90)
            result['files'].append(out)
            result['stdout'] = f'Saved: {out}'

        # ── Image → PDF ──
        elif from_format in ('png', 'jpg', 'jpeg', 'webp', 'bmp', 'gif', 'tiff') and to_format == 'pdf':
            from PIL import Image
            paths = [input_path]
            images = []
            for p in paths:
                img = Image.open(p).convert('RGB')
                images.append(img)
            out = _outpath(input_path, 'pdf', output_dir, 'output')
            if len(images) == 1:
                images[0].save(out, 'PDF')
            else:
                images[0].save(out, 'PDF', save_all=True, append_images=images[1:])
            result['files'].append(out)
            result['stdout'] = f'Saved: {out}'

        # ── Text/Document → PDF ──
        elif from_format in ('txt', 'md', 'html') and to_format == 'pdf':
            try:
                from fpdf import FPDF
            except ImportError:
                raise ValueError(
                    'Document → PDF conversion requires fpdf2. '
                    'Install fpdf2 and try again.'
                )
            text = _read_file(input_path, input_text)
            pdf = FPDF()
            pdf.add_page()
            pdf.add_font('NotoSansSC', '', os.path.join(os.path.dirname(__file__), 'NotoSansSC-Regular.ttf'),
                         uni=True) if os.path.exists(os.path.join(os.path.dirname(__file__),
                                                                  'NotoSansSC-Regular.ttf')) else None
            pdf.set_auto_page_break(auto=True, margin=15)

            if from_format == 'html':
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(text, 'html.parser')
                text = soup.get_text()

            # Try to use a Unicode font; fall back to default
            try:
                pdf.add_font('CJK', '', '', uni=True)
                pdf.set_font('CJK', size=11)
            except:
                try:
                    pdf.set_font('Helvetica', size=11)
                except:
                    pdf.set_font('Times', size=11)

            for line in text.split('\n'):
                s = line.strip()
                if not s:
                    pdf.ln(5)
                    continue
                # Handle markdown headings
                if from_format == 'md' and s.startswith('#'):
                    level = min(len(s.split(' ')[0]), 4)
                    size = [24, 18, 14, 12][level - 1]
                    try:
                        pdf.set_font('CJK', size=size)
                    except:
                        pdf.set_font('Helvetica', size=size)
                    pdf.multi_cell(0, 10, s.lstrip('#').strip())
                    try:
                        pdf.set_font('CJK', size=11)
                    except:
                        pdf.set_font('Helvetica', size=11)
                else:
                    pdf.multi_cell(0, 7, s)

            out = _outpath(input_path, 'pdf', output_dir)
            pdf.output(out)
            result['files'].append(out)
            result['stdout'] = f'Saved: {out}'

        # ── EPUB extraction ──
        elif from_format == 'epub' and to_format in ('txt', 'md'):
            try:
                import ebooklib
                from ebooklib import epub
            except ImportError:
                raise ValueError(
                    'EPUB conversion requires ebooklib. '
                    'Install ebooklib and try again.'
                )
            book = epub.read_epub(input_path)
            lines = []
            title = book.get_metadata('DC', 'title')
            if title:
                if to_format == 'md':
                    lines.append(f'# {title[0][0]}\n')
                else:
                    lines.append(f'{title[0][0]}\n{"=" * len(title[0][0])}\n')

            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    content = item.get_body_content()
                    if content:
                        from bs4 import BeautifulSoup
                        soup = BeautifulSoup(content, 'html.parser')
                        text = soup.get_text(strip=True)
                        if text:
                            lines.append(text)
            result['stdout'] = '\n\n'.join(lines)

        # ── Image format conversion (extended: bmp/gif/tiff support via Pillow) ──
        _IMG_FMTS = {'png','jpg','jpeg','webp','bmp','gif','tiff'}
        if from_format in _IMG_FMTS and to_format in _IMG_FMTS:
            from PIL import Image
            img = Image.open(input_path)
            out = _outpath(input_path, to_format, output_dir)
            fmt_map = {'png':'PNG','jpg':'JPEG','jpeg':'JPEG','webp':'WEBP',
                       'bmp':'BMP','gif':'GIF','tiff':'TIFF'}
            fmt = fmt_map.get(to_format, to_format.upper())
            quality = 90 if fmt in ('JPEG', 'WEBP') else None
            if fmt == 'GIF':
                img = img.convert('P', palette=Image.Palette.ADAPTIVE)
            elif fmt == 'JPEG' and img.mode in ('RGBA','P'):
                bg = Image.new('RGB', img.size, (255,255,255))
                bg.paste(img, mask=img.split()[-1] if img.mode=='RGBA' else None)
                img = bg
            img.save(out, fmt, quality=quality)
            result['files'].append(out)
            result['stdout'] = f'Saved: {out}'

        else:
            raise ValueError(f'Conversion from {from_format} to {to_format} not supported')

    except Exception as e:
        result['error'] = f'{type(e).__name__}: {str(e)}'
        result['stdout'] = ''

    return json.dumps(result)


def _read_file(path, text):
    if path:
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    return text or ''


def _outpath(input_path, ext, output_dir, fallback_name='output'):
    if input_path:
        base = os.path.basename(input_path).rsplit('.', 1)[0]
    else:
        base = fallback_name
    return os.path.join(output_dir, f'{base}.{ext}')


if __name__ == '__main__':
    input_path = sys.argv[1] if len(sys.argv) > 1 else ''
    input_text = sys.argv[2] if len(sys.argv) > 2 else ''
    from_format = sys.argv[3] if len(sys.argv) > 3 else 'txt'
    to_format = sys.argv[4] if len(sys.argv) > 4 else 'txt'
    output_dir = sys.argv[5] if len(sys.argv) > 5 else '/storage/emulated/0/Download'
    print(convert(input_path, input_text, from_format, to_format, output_dir))
